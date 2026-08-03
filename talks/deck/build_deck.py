#!/usr/bin/env python3
"""Build the vector-search talk deck (.pptx) — imports cleanly into Google Slides."""
import json, os
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs")
N = json.load(open(os.path.join(FIGS, "numbers.json")))

# ── palette ──────────────────────────────────────────────────────────────────
BG      = C(0x10, 0x13, 0x1a)
PANEL   = C(0x17, 0x1b, 0x24)
FG      = C(0xe8, 0xed, 0xf5)
MUTED   = C(0x93, 0xa1, 0xb5)
DIM     = C(0x5b, 0x66, 0x78)
BLUE    = C(0x4d, 0xa3, 0xff)
AMBER   = C(0xff, 0xb4, 0x54)
GREEN   = C(0x4a, 0xde, 0x80)
RED     = C(0xf8, 0x71, 0x71)
PINK    = C(0xff, 0x2d, 0x55)
LINE    = C(0x2b, 0x33, 0x42)

FONT = "Arial"
MONO = "Courier New"

prs = Presentation()
prs.slide_width, prs.slide_height = In(13.333), In(7.5)
W, H = 13.333, 7.5
BLANK = prs.slide_layouts[6]

deck_notes = []


def slide(note=""):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background(); bg.shadow.inherit = False
    if note:
        s.notes_slide.notes_text_frame.text = note
    deck_notes.append(note)
    return s


def tx(s, text, l, t, w, h, size=20, color=FG, bold=False, align=PP_ALIGN.LEFT,
       font=FONT, italic=False, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return box


def rule(s, l, t, w, color=BLUE, h=0.045):
    sh = s.shapes.add_shape(1, In(l), In(t), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def panel(s, l, t, w, h, color=PANEL):
    sh = s.shapes.add_shape(5, In(l), In(t), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.color.rgb = LINE; sh.line.width = Pt(1); sh.shadow.inherit = False
    return sh


def head(s, title, kicker=None):
    if kicker:
        tx(s, kicker.upper(), 0.85, 0.52, 11, 0.3, 13, BLUE, True, spacing=1)
        tx(s, title, 0.85, 0.85, 11.6, 0.9, 32, FG, True)
        rule(s, 0.85, 1.72, 1.5)
    else:
        tx(s, title, 0.85, 0.62, 11.6, 0.9, 32, FG, True)
        rule(s, 0.85, 1.52, 1.5)


def bullets(s, items, top=2.1, size=19, gap=0.62, left=0.85, width=11.6):
    y = top
    for it in items:
        if isinstance(it, tuple):
            text, col, bold = (it + (False,))[:3] if len(it) < 3 else it
        else:
            text, col, bold = it, FG, False
        if text == "":
            y += gap * 0.45; continue
        dot = s.shapes.add_shape(9, In(left), In(y + 0.13), In(0.13), In(0.13))
        dot.fill.solid(); dot.fill.fore_color.rgb = col if col != FG else BLUE
        dot.line.fill.background(); dot.shadow.inherit = False
        tx(s, text, left + 0.38, y, width - 0.38, gap, size, col, bold)
        y += gap + 0.12 * text.count("\n")
    return y


def table(s, headers, rows, col_w, top=2.2, left=0.85, size=15, hi=None,
          hi_color=AMBER, row_h=0.44):
    tot = sum(col_w)
    shape = s.shapes.add_table(len(rows) + 1, len(headers),
                               In(left), In(top), In(tot), In(row_h * (len(rows) + 1)))
    t = shape.table
    t.first_row = False; t.horz_banding = False
    for i, w in enumerate(col_w):
        t.columns[i].width = In(w)
    for r in range(len(rows) + 1):
        t.rows[r].height = In(row_h)
    for c, htxt in enumerate(headers):
        cell = t.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = BG
        cell.margin_left = cell.margin_right = In(0.09)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(size - 2); r.font.bold = True; r.font.color.rgb = MUTED; r.font.name = FONT
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if ri == hi else BG
            cell.margin_left = cell.margin_right = In(0.09)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(size)
            r.font.bold = (ri == hi)
            r.font.color.rgb = hi_color if ri == hi else FG
            r.font.name = FONT
    return shape


def image(s, name, top=1.95, height=4.9, note=None):
    path = os.path.join(FIGS, name)
    from PIL import Image
    iw, ih = Image.open(path).size
    h = height; w = h * iw / ih
    if w > 11.6:
        w = 11.6; h = w * ih / iw
    s.shapes.add_picture(path, In((W - w) / 2), In(top), In(w), In(h))
    if note:
        tx(s, note, 0.85, top + h + 0.16, 11.6, 0.5, 15, MUTED, align=PP_ALIGN.CENTER)


def stat(s, value, caption, sub=None, color=AMBER, top=2.35, vsize=112):
    """Big number + optional caption + optional sub-line, stacked adaptively."""
    vh = vsize / 72.0 * 1.32                       # actual height of the numeral line
    tx(s, value, 0.85, top, 11.6, vh, vsize, color, True, PP_ALIGN.CENTER)
    y = top + vh + 0.08
    if caption:
        tx(s, caption, 0.85, y, 11.6, 0.5, 26, FG, True, PP_ALIGN.CENTER)
        y += 0.62
    if sub:
        tx(s, sub, 1.0, y, 11.3, 0.8, 17, MUTED, align=PP_ALIGN.CENTER)


def section(num, title, sub, note=""):
    s = slide(note)
    tx(s, num, 0.9, 2.3, 2.4, 1.8, 128, BLUE, True)
    tx(s, title, 3.3, 2.55, 9.2, 1.2, 50, FG, True)
    rule(s, 3.35, 3.95, 2.2, AMBER)
    tx(s, sub, 3.35, 4.3, 8.9, 1.6, 21, MUTED)
    return s


def code(s, lines, top=2.2, left=1.0, width=11.3, size=17, height=None):
    h = height or (0.36 * len(lines) + 0.5)
    panel(s, left, top, width, h, C(0x0b, 0x0e, 0x14))
    y = top + 0.25
    for ln in lines:
        col = MUTED if ln.strip().startswith("#") else FG
        tx(s, ln if ln else " ", left + 0.35, y, width - 0.7, 0.36, size, col, font=MONO)
        y += 0.36
    return top + h


# ═════════════════════════════════════════════════════════════════════ 1 TITLE
s = slide("Welcome. Sixty minutes, four techniques. Everything I claim today, I measured — "
          "there's a notebook and every number comes out of it.")
tx(s, "VECTOR SEARCH", 0.9, 2.15, 11.6, 0.5, 17, BLUE, True, spacing=1)
tx(s, "From brute force to HNSW", 0.9, 2.6, 11.6, 1.2, 54, FG, True)
rule(s, 0.92, 4.0, 2.6, AMBER)
tx(s, "kNN  ·  ANN  ·  IVF-PQ  ·  HNSW", 0.9, 4.4, 11.6, 0.6, 26, MUTED)
tx(s, "Every number in this deck is measured, not estimated.\n"
      "Companion notebook: notebooks/talk-vector-search.ipynb", 0.9, 5.6, 11.6, 1.0, 15, DIM)

# ══════════════════════════════════════════════════════════════════ 2 OPENING
s = slide("Embeddings turn anything into a list of numbers, and similar things land near each "
          "other. So 'find me things like this' becomes a geometry problem.")
head(s, "Everything is a vector now")
bullets(s, [
    ("A document, an image, a user, a support ticket → a list of 256–1536 numbers.", FG),
    ("Similar things land near each other in that space.", FG),
    ("", FG),
    ('So "find me things like this" stops being a text problem…', MUTED),
    ("…and becomes a geometry problem: given a query vector, find the closest vectors.", GREEN, True),
], top=2.3, gap=0.7)

# ═════════════════════════════════════════════════════════════════ 3 AGENDA
s = slide("Four techniques, in the order they were invented. Each one exists because of a "
          "problem with the one before it. That's the whole talk.")
head(s, "Four techniques. That's the whole talk.")
rows = [
    ("1.  kNN", "Exact search. Correct, simple, stops scaling."),
    ("2.  ANN", "Trade a little correctness for a lot of speed."),
    ("3.  IVF-PQ", "Partition the space, then compress what's in it."),
    ("4.  HNSW", "Build a graph and walk it."),
]
y = 2.25
for i, (k, v) in enumerate(rows):
    panel(s, 0.85, y, 11.6, 0.95)
    tx(s, k, 1.15, y + 0.24, 2.7, 0.5, 24, BLUE, True)
    tx(s, v, 4.0, y + 0.28, 8.2, 0.5, 19, FG)
    y += 1.12
tx(s, "Each one exists because of a problem with the one before it.",
   0.85, y + 0.12, 11.6, 0.5, 17, MUTED, italic=True)

# ═══════════════════════════════════════════════════════════════ 4 SUPPORTBOT
s = slide("One problem carries the whole talk so every technique is scored on the same task. "
          "Swap in whatever your team works on — the numbers barely change.")
head(s, "One problem carries the whole talk", "the running example")
panel(s, 0.85, 2.0, 11.6, 2.15)
tx(s, "SupportBot", 1.25, 2.25, 10.8, 0.5, 23, AMBER, True)
tx(s, "200,000 past support tickets. A new ticket arrives — surface the most similar\n"
      "past tickets so the agent can reuse the resolution.", 1.25, 2.8, 10.8, 1.1, 19, FG, spacing=1.25)
y = 4.45
for i, (label, val, col) in enumerate([("vectors", f"{N['N_TICKETS']:,}", FG),
                        ("dimensions", f"{N['N_DIMS']}", FG),
                        ("throughput", "200 QPS", FG),
                        ("raw size", f"{N['raw_bytes']/1e6:.0f} MB", AMBER)]):
    x = 0.85 + 2.95 * i
    panel(s, x, y, 2.75, 1.35)
    tx(s, label.upper(), x + 0.25, y + 0.22, 2.3, 0.3, 12, MUTED, True)
    tx(s, val, x + 0.25, y + 0.58, 2.3, 0.6, 27, col, True)
tx(s, "Remember 205 MB — in about twenty-five minutes I'm going to make it six.",
   0.85, 6.3, 11.6, 0.5, 17, MUTED, italic=True)

# ═════════════════════════════════════════════════════════════════ 5 METRICS
s = slide("The thing nobody tells you up front: normalize once at write time and all three "
          "metrics rank identically. Then stop thinking about it.")
head(s, "One piece of vocabulary first")
bullets(s, [
    ("Three distance metrics you'll see: Euclidean, dot product, cosine.", FG),
    ("", FG),
    ("If you normalize vectors to length 1 at write time, all three rank IDENTICALLY.", GREEN, True),
    ("", FG),
    ("So normalize once, then stop thinking about it.", MUTED),
], top=2.4, gap=0.72)

# ══════════════════════════════════════════════════════════ SECTION 1: kNN
section("1", "Exact kNN", "Compare the query to every vector. Keep the best k.\n"
        "It's a full scan with a heap — and it's better than you think.",
        "Let's start with the dumbest thing that works, because it's better than you think.")

s = slide("That is a complete, production-viable search engine. One matrix-vector product, one "
          "partial sort. If your team is standing up a vector DB for 50k rows, this slide is the talk.")
head(s, "The entire search engine", "1 · exact knn")
code(s, [
    "scores = tickets @ query                 # cosine sim to every ticket",
    "top    = np.argpartition(-scores, k)[:k] # k best, unordered",
    "return top[np.argsort(-scores[top])]     # order just those k",
], top=2.2, size=18)
tx(s, "If your team is standing up a vector database for 50,000 rows —\nthis slide is the whole talk.",
   0.85, 4.4, 11.6, 1.1, 26, AMBER, True, spacing=1.2)

s = slide("Two values, and the second is the one people forget. It's your ruler: every "
          "approximate method gets measured against it, forever.")
head(s, "What exact kNN buys you", "1 · exact knn")
panel(s, 0.85, 2.05, 5.6, 3.9)
tx(s, "VALUE 1", 1.2, 2.3, 4.9, 0.3, 13, BLUE, True)
tx(s, "Simplest thing\nthat works", 1.2, 2.65, 4.9, 1.0, 26, FG, True, spacing=1.15)
tx(s, "No index. No training. No tuning.\nNo new service.\n\nAdding a vector is append().\nDeleting is a mask.",
   1.2, 3.72, 4.9, 1.7, 15, MUTED, spacing=1.32)
panel(s, 6.85, 2.05, 5.6, 3.9)
tx(s, "VALUE 2", 7.2, 2.3, 4.9, 0.3, 13, AMBER, True)
tx(s, "It is your ruler", 7.2, 2.65, 4.9, 1.0, 26, FG, True, spacing=1.15)
tx(s, "Every approximate method is\nscored against it.\n\nRun it once on 1,000 queries —\nthat cached file is your only alarm.",
   7.2, 3.72, 4.9, 1.7, 15, MUTED, spacing=1.32)
tx(s, "You keep exact search in your stack forever — even after you stop serving from it.",
   0.85, 6.25, 11.6, 0.5, 18, GREEN, True, PP_ALIGN.CENTER)

s = slide("[Pause — this is deliberately anticlimactic.] Which is completely fine. At 200,000 "
          "vectors you do not need any of the next fifty minutes.")
head(s, "So where does it fall over?", "1 · exact knn")
cores = 200 * N["EXACT_MS"] / 1000
stat(s, f"{N['EXACT_MS']:.1f} ms", "per query, on 200,000 tickets",
     f"= {cores:.1f} of a CPU core to sustain 200 QPS", AMBER, top=2.1, vsize=108)
tx(s, "…which is completely fine. Nothing is wrong.\n"
      "At 200,000 vectors you do not need any of the next fifty minutes.",
   0.85, 5.85, 11.6, 1.0, 20, GREEN, True, PP_ALIGN.CENTER, spacing=1.25)

s = slide("Here's where it breaks. The cost is exactly linear — ten times the tickets, ten times "
          "the bill. And you blew the 50ms p99 somewhere around two million.")
head(s, "It doesn't fail. It just gets linearly more expensive.", "1 · exact knn")
ms = N["EXACT_MS"]
table(s, ["corpus", "per query", "cores @ 200 QPS", ""],
      [(f"{N['N_TICKETS']:,}", f"{ms:.1f} ms", f"{cores:.1f}", "fine"),
       (f"{N['N_TICKETS']*10:,}", f"{ms*10:.0f} ms", f"{cores*10:.1f}", "p99 blown"),
       (f"{N['N_TICKETS']*100:,}", f"{ms*100:.0f} ms",
        f"{cores*100:.1f}", "a rack of machines")],
      [3.0, 2.6, 3.0, 3.0], top=2.35, size=19, hi=2, row_h=0.66)
tx(s, "Ten times the tickets, ten times the bill. Forever.",
   0.85, 5.0, 11.6, 0.5, 21, AMBER, True)
tx(s, "Seventy cores, burning continuously, to answer a question that ought to fit on a laptop.",
   0.85, 5.6, 11.6, 0.5, 17, MUTED)

s = slide("The obvious question: isn't this solved? kd-trees work in 2-3 dimensions and are "
          "useless at 256. There is no exact escape hatch — that's why this field exists.")
head(s, "Isn't this solved? kd-trees, R-trees?", "1 · exact knn")
bullets(s, [
    ("They genuinely work in 2 or 3 dimensions.", FG),
    ("They are useless at 256 — above ~20 dims they degrade to scanning everything, plus overhead.", RED),
    ("", FG),
    ("THE CURSE OF DIMENSIONALITY", AMBER, True),
    ("In high dimensions almost every point sits at almost the same distance from your query.", MUTED),
    ("Distances bunch up — and if everything is equidistant, there is no region to prune.", MUTED),
], top=2.25, gap=0.62)
tx(s, "There is no exact escape hatch. That's not a gap in the literature —\nit's the reason this entire field exists.",
   0.85, 6.0, 11.6, 1.0, 20, GREEN, True, spacing=1.2)

# ══════════════════════════════════════════════════════════ SECTION 2: ANN
section("2", "ANN", "Trade a little correctness for a lot of speed.\n"
        "Then measure exactly how much you traded.",
        "What we give up is being exactly right, and we give up surprisingly little of it.")

s = slide("Recall@k is the whole definition. You asked for ten, you got nine right ones, "
          "recall@10 is 0.9. You compute it by diffing against exact search.")
head(s, "recall@k — the one metric", "2 · ann")
panel(s, 0.85, 2.05, 11.6, 1.5)
tx(s, "The fraction of the true top-k that you actually returned.",
   1.2, 2.45, 10.9, 0.7, 25, FG, True, PP_ALIGN.CENTER)
tx(s, "You asked for 10.  You got 9 of the right ones.  →  recall@10 = 0.90",
   0.85, 3.85, 11.6, 0.6, 22, AMBER, True, PP_ALIGN.CENTER)
bullets(s, [
    ("Compute it by running exact search on ~1,000 sampled queries and diffing.", MUTED),
    ("That cached ground truth is the ruler from section 1.", MUTED),
], top=4.75, gap=0.6)

s = slide("This is the objection you should get. The answer: your embeddings were already "
          "approximate. But name the exception honestly — dedup is a different regime.")
head(s, '"Isn\'t 0.9 just… accepting wrong answers?"', "2 · ann")
panel(s, 0.85, 2.05, 5.6, 3.6)
tx(s, "WHY IT'S FINE", 1.2, 2.3, 4.9, 0.3, 13, GREEN, True)
tx(s, "Your embeddings were\nalready approximate.", 1.2, 2.7, 4.9, 1.0, 22, FG, True, spacing=1.15)
tx(s, "The true #7 by cosine distance\nis not a ground truth about\nrelevance. It's one model's opinion.\n\nChasing exact geometry on a\nfuzzy measurement is false precision.",
   1.2, 3.82, 4.9, 1.7, 14.5, MUTED, spacing=1.3)
panel(s, 6.85, 2.05, 5.6, 3.6)
tx(s, "WHERE IT ISN'T", 7.2, 2.3, 4.9, 0.3, 13, RED, True)
tx(s, "Dedup and exact\nretrieval.", 7.2, 2.7, 4.9, 1.0, 22, FG, True, spacing=1.15)
tx(s, "There a missed neighbour isn't a\nranking nudge — it's a\ncorrectness bug.\n\nRecall 0.95 means 5% wrong\nanswers. Different conversation.",
   7.2, 3.82, 4.9, 1.7, 14.5, MUTED, spacing=1.3)
tx(s, "First question about any vector search feature: which regime am I in?",
   0.85, 5.95, 11.6, 0.5, 19, AMBER, True, PP_ALIGN.CENTER)

s = slide("Everything from here is one of three ideas. The first two combine, which is why "
          "IVF-PQ is one word.")
head(s, "Three ideas. That's all that's left.", "2 · ann")
for i, (t, d, col) in enumerate([
        ("Partition", "Scan a slice of the data instead of all of it.", BLUE),
        ("Compress", "Make each comparison cheaper.", AMBER),
        ("Build a graph", "Walk toward the answer in log time.", GREEN)]):
    x = 0.85 + i * 3.95
    panel(s, x, 2.3, 3.65, 2.6)
    tx(s, t, x + 0.35, 2.65, 3.0, 0.6, 25, col, True)
    tx(s, d, x + 0.35, 3.4, 3.0, 1.3, 17, MUTED, spacing=1.3)
tx(s, "The first two combine — which is why you'll only ever see them written as one word: IVF-PQ.",
   0.85, 5.4, 11.6, 0.5, 19, FG, True, PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════ SECTION 3: IVF-PQ
section("3", "IVF-PQ", "Two independent ideas bolted together.\n"
        "IVF partitions (buys time) · PQ compresses (buys memory).",
        "It has a compound name because it is literally two ideas. I'll teach them separately, "
        "measure each alone, then combine.")

s = slide("The library analogy. You don't scan every book in the building — you walk to the "
          "right few shelves and scan those.")
head(s, "3a · IVF — Inverted File Index", "partition · buys time")
panel(s, 0.85, 2.0, 11.6, 1.15)
tx(s, "You don't scan every book in the building. You walk to the right few shelves.",
   1.2, 2.35, 10.9, 0.6, 23, AMBER, True, PP_ALIGN.CENTER)
bullets(s, [
    ("BUILD  —  k-means the corpus into 512 buckets (Voronoi cells). ~390 tickets each.", FG),
    ("SEARCH  —  compare the query to the 512 centroids only. That's cheap.\nPick the nprobe closest buckets. Brute-force just those.", FG),
], top=3.5, gap=0.72)
tx(s, "Value: you stop looking at data that was never going to win.",
   0.85, 5.85, 11.6, 0.5, 21, GREEN, True)

ivf = {r["nprobe"]: r for r in N["ivf_sweep"]}
s = slide("Two things to point at. One row here is a product decision. And recall climbs fast "
          "then flattens, while cost keeps climbing linearly and never flattens.")
head(s, "The nprobe sweep — the dial", "3a · ivf")
rows = []
for p in [1, 2, 4, 8, 16, 64, 512]:
    r = ivf[p]
    rows.append((f"{p}", f"{r['scanned']:,.0f}", f"{100*r['scanned']/N['N_TICKETS']:.1f}%",
                 f"{r['ms']:.3f} ms", f"{r['recall']:.3f}"))
table(s, ["nprobe", "scanned", "% of corpus", "ms/query", "recall@10"], rows,
      [1.9, 2.4, 2.4, 2.4, 2.5], top=2.1, size=16, hi=2, row_h=0.5)
tx(s, f"At nprobe=4 we scanned 0.8% of the corpus and got 96% of the right answers — "
      f"a ~95× speedup for 4 points of recall.", 0.85, 6.15, 11.6, 0.6, 19, AMBER, True)

s = slide("Almost all of the value is in the first few probes. 1→4 buys 37 points of recall. "
          "64→512 buys zero and costs 3.2 milliseconds.")
head(s, "Recall saturates. Cost doesn't.", "3a · ivf")
image(s, "ivf_sweep.png", top=1.95, height=4.0)
tx(s, "1 → 4 probes:  +37 points of recall.        64 → 512 probes:  +0.000, and 3.2 ms slower.",
   0.85, 6.15, 11.6, 0.5, 19, AMBER, True, PP_ALIGN.CENTER)

s = slide("This is IVF's whole failure mode, and I love it because it's specific and "
          "predictable. Errors you can explain are errors you can debug.")
head(s, "Where IVF is wrong — and why that's good", "3a · ivf")
image(s, "voronoi.png", top=1.9, height=4.35)
tx(s, "The query hugged a boundary. Its true nearest neighbour fell one cell over. nprobe is exactly the fix.",
   0.85, 6.4, 11.6, 0.5, 17, MUTED, align=PP_ALIGN.CENTER)

s = slide("Three things. And hold the memory one — it's the setup for PQ in ninety seconds.")
head(s, "Where IVF's value runs out", "3a · ivf")
bullets(s, [
    ("Memory is completely unchanged — still all 205 MB of full float32 vectors.", RED, True),
    ("IVF buys TIME, not SPACE. Hold that thought for ninety seconds.", MUTED),
    ("", FG),
    ("It needs training — k-means over a sample before you can insert anything.", FG),
    ("", FG),
    ("It drifts. Centroids fitted on last year's tickets go lopsided as topics shift.", FG),
    ("Retraining is real, scheduled, operational work.", MUTED),
], top=2.2, gap=0.58)

# ── 3b PQ ───────────────────────────────────────────────────────────────────
s = slide("The palette analogy: instead of storing every pixel's exact colour, store an index "
          "into a 256-colour palette.")
head(s, "3b · PQ — Product Quantization", "compress · buys memory")
panel(s, 0.85, 2.0, 11.6, 1.15)
tx(s, "Don't store every pixel's exact colour. Store an index into a 256-colour palette.",
   1.2, 2.35, 10.9, 0.6, 22, AMBER, True, PP_ALIGN.CENTER)
bullets(s, [
    ("SPLIT  —  chop each 256-dim vector into 32 sub-vectors of 8 dims.", FG),
    ("CLUSTER  —  k-means each chunk position independently → 32 codebooks of 256 centroids.", FG),
    ("ENCODE  —  replace each chunk with its nearest centroid id: 0–255. One byte.", FG),
], top=3.5, gap=0.72)
tx(s, "1,024 bytes  →  32 bytes", 0.85, 5.85, 11.6, 0.7, 30, GREEN, True, PP_ALIGN.CENTER)

s = slide("There it is. Ticket zero used to be 256 floats. It is now literally these 32 bytes.")
head(s, "A ticket, before and after", "3b · pq")
panel(s, 0.85, 2.1, 11.6, 1.35)
tx(s, "BEFORE", 1.2, 2.32, 2.0, 0.3, 13, MUTED, True)
tx(s, "[-0.0872  0.0147  0.0100  -0.0640  -0.0553  … ]   256 floats  =  1,024 bytes",
   1.2, 2.72, 10.9, 0.5, 17, FG, font=MONO)
panel(s, 0.85, 3.65, 11.6, 1.35)
tx(s, "AFTER", 1.2, 3.87, 2.0, 0.3, 13, GREEN, True)
tx(s, "[125 195 131  51 129  60  73  33 … ]   32 bytes",
   1.2, 4.27, 10.9, 0.5, 17, GREEN, font=MONO)
tx(s, f"{N['raw_bytes']/1e6:.0f} MB  →  {N['pq_bytes']/1e6:.1f} MB", 0.85, 5.3, 11.6, 0.75, 40, AMBER, True, PP_ALIGN.CENTER)
tx(s, f"{N['raw_bytes']/N['pq_bytes']:.0f}× smaller", 0.85, 6.15, 11.6, 0.5, 23, FG, True, PP_ALIGN.CENTER)

s = slide("This is a deployment value, not a latency value. Three arguments live in this table "
          "and they get better as you go down.")
head(s, "PQ's value is a DEPLOYMENT value", "3b · pq")
d, m = N["N_DIMS"], 32
table(s, ["corpus", "raw", "with PQ", "what changes"],
      [(f"{N['N_TICKETS']:,}", f"{N['N_TICKETS']*d*4/1e9:.2f} GB", f"{N['N_TICKETS']*m/1e9:.2f} GB",
        "fits inside your app process"),
       (f"{N['N_TICKETS']*10:,}", f"{N['N_TICKETS']*10*d*4/1e9:.1f} GB", f"{N['N_TICKETS']*10*m/1e9:.2f} GB",
        "growth stopped being an infra project"),
       ("100,000,000", "102 GB", "3.2 GB", "one big machine, not a cluster"),
       ("1,000,000,000", "1,024 GB", "32 GB", "billion-scale on a single host")],
      [2.5, 1.9, 1.8, 5.4], top=2.3, size=16, hi=3, row_h=0.62)
tx(s, "The top row deletes a service from your architecture diagram. The bottom row is why PQ exists.",
   0.85, 5.5, 11.6, 0.5, 18, AMBER, True)

pqs = {r["m"]: r for r in N["pq_sweep"]}
s = slide("Read these and don't panic — they look alarming on purpose. At 32x compression the "
          "reconstruction error is comparable to the gap between the true #1 and the true #40.")
head(s, "But the compression dial has a CLIFF", "3b · pq")
table(s, ["m", "bytes/vector", "compression", "recall@10"],
      [(f"{m}", f"{m}", f"{pqs[m]['ratio']:.0f}×", f"{pqs[m]['recall']:.3f}") for m in [16, 32, 64, 128]],
      [1.8, 2.8, 2.8, 2.8], top=2.35, size=17, hi=1, row_h=0.58)
tx(s, "PQ alone is not a search index. It's a SHORTLISTING index.",
   0.85, 5.35, 11.6, 0.7, 26, RED, True, PP_ALIGN.CENTER)
tx(s, "Hold that for five minutes — three lines of code fix almost all of it.",
   0.85, 6.1, 11.6, 0.5, 18, MUTED, align=PP_ALIGN.CENTER)

s = slide("[Slow down. This is the best idea in the talk.] The natural assumption is that "
          "compressed data must be decompressed. PQ never does. Not once.")
head(s, "The clever part: you never decompress", "3b · pq")
bullets(s, [
    ("Squared distance is SEPARABLE — the distance between two vectors is just the sum\nof the distances between their chunks. That's not a trick, it's arithmetic.", MUTED),
], top=2.05, gap=0.5)
panel(s, 0.85, 3.15, 11.6, 1.5)
tx(s, "Precompute one lookup table per query:\ndistance from each query chunk to each of its 256 centroids.  →  32 × 256 = 8,192 tiny computations, ONCE.",
   1.2, 3.42, 10.9, 1.0, 18, FG, spacing=1.35)
tx(s, "Then the distance to ANY ticket = 32 table lookups + 32 adds.",
   0.85, 5.0, 11.6, 0.7, 27, GREEN, True, PP_ALIGN.CENTER)
tx(s, "No multiplications. No 256-dimensional math. The same 32 bytes that saved memory\nmade the comparison cheaper too — compression and speed usually trade against each other.",
   0.85, 5.85, 11.6, 1.0, 17, MUTED, align=PP_ALIGN.CENTER, spacing=1.3)

s = slide("Right on average, noisy in the details. The scatter around the line IS the recall "
          "loss — and the top-10 live in the bottom-left corner, where the noise hurts.")
head(s, "Right on average. Noisy where it matters.", "3b · pq")
image(s, "pq_scatter.png", top=1.9, height=4.3)
tx(s, "The scatter around the line is the recall loss. The top-10 live in the bottom-left — exactly where the noise bites.",
   0.85, 6.35, 11.6, 0.5, 17, MUTED, align=PP_ALIGN.CENTER)

# ── 3c combined + rerank ────────────────────────────────────────────────────
s = slide("Now put them together. IVF picks the buckets, PQ makes scanning them cheap. "
          "It's fast, it's tiny, and its recall is 0.485 — useless.")
head(s, "3c · IVF-PQ, combined", "partition + compress")
panel(s, 0.85, 2.0, 11.6, 1.15)
tx(s, '"Go to the right shelves — and read compressed summaries once you\'re there."',
   1.2, 2.35, 10.9, 0.6, 22, AMBER, True, PP_ALIGN.CENTER)
for i, (lbl, val, col) in enumerate([
        ("latency", f"{N['MS_IVFPQ']:.3f} ms", GREEN),
        ("memory", f"{N['pq_bytes']/1e6:.1f} MB", GREEN),
        ("recall@10", f"{N['REC_IVFPQ']:.3f}", RED)]):
    x = 0.85 + i * 3.95
    panel(s, x, 3.55, 3.65, 1.75)
    tx(s, lbl.upper(), x + 0.35, 3.8, 3.0, 0.3, 13, MUTED, True)
    tx(s, val, x + 0.35, 4.2, 3.0, 0.7, 32, col, True)
tx(s, "So we're fast, we're small, and we're wrong. Let's fix the wrong part.",
   0.85, 5.7, 11.6, 0.6, 22, FG, True, PP_ALIGN.CENTER)

s = slide("Retrieve top-100 by cheap PQ distance. Load those 100 full vectors. Rescore exactly. "
          "A hundred exact distance computations is nothing.")
head(s, "Rerank — three lines", "3c · ivf-pq")
code(s, [
    "_, candidates = ivfpq.search(Q, 100)    # cheap, approximate shortlist",
    "exact = tickets[candidates] @ query     # exact, on 100 vectors only",
    "return candidates[argsort(-exact)[:10]] # true top-10 of the shortlist",
], top=2.15, size=17)
tx(s, "100 exact distance computations is nothing.", 0.85, 4.25, 11.6, 0.5, 20, MUTED)
stat(s, f"{N['REC_IVFPQ']:.3f}  →  {N['REC_RERANK']:.3f}", "recall@10",
     f"for +{(N['MS_RERANK']-N['MS_IVFPQ'])*1000:.0f} microseconds", GREEN, top=4.75, vsize=52)

s = slide("If this room takes home one implementation detail, make it this one. It is routinely "
          "missing from hand-rolled vector search.")
head(s, "The highest payoff in the whole stack", "3c · ivf-pq")
table(s, ["candidates", "ms/query", "recall@10", "gain"],
      [("— (no rerank)", f"{N['MS_IVFPQ']:.3f}", f"{N['REC_IVFPQ']:.3f}", "—"),
       ("20", "0.100", "0.691", "+0.206"),
       ("50", "0.112", "0.903", "+0.418"),
       ("100", f"{N['MS_RERANK']:.3f}", f"{N['REC_RERANK']:.3f}",
        f"+{N['REC_RERANK']-N['REC_IVFPQ']:.3f}"),
       ("200", "0.161", "0.996", "+0.512")],
      [3.0, 2.7, 2.8, 2.7], top=2.3, size=17, hi=3, row_h=0.55)
tx(s, "Every few months a team concludes “we tried PQ, quality was terrible.”\n"
      "What they shipped was a shortlisting index without the second stage.",
   0.85, 5.85, 11.6, 1.0, 18, AMBER, True, spacing=1.25)

s = slide("Name the cost of your own fix — that's what makes the rest of your numbers "
          "believable. Reranking needs the vectors you just compressed away.")
head(s, "The honest catch", "3c · ivf-pq")
panel(s, 0.85, 2.05, 11.6, 1.0)
tx(s, "Reranking needs the full vectors — which you just spent a section compressing away.",
   1.2, 2.32, 10.9, 0.5, 21, RED, True, PP_ALIGN.CENTER)
y = 3.35
for lbl, txt, col in [
        ("On SSD", "~100 vectors per query, a few hundred KB. Keeps the memory win. Usually right.", GREEN),
        ("In RAM", "Fastest — but you just gave back the entire 32× saving. Sometimes correct; know you did it.", AMBER)]:
    panel(s, 0.85, y, 11.6, 1.2)
    tx(s, lbl, 1.2, y + 0.35, 2.1, 0.5, 21, col, True)
    tx(s, txt, 3.5, y + 0.4, 8.6, 0.6, 17, FG)
    y += 1.4
tx(s, "In practice it's SSD, and you keep the 32× win that made you choose IVF-PQ in the first place.",
   0.85, 6.25, 11.6, 0.5, 18, MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════ SECTION 4: HNSW
section("4", "HNSW", "Hierarchical Navigable Small World.\n"
        "Don't partition. Don't compress. Build a graph and walk it.",
        "Completely different philosophy. Terrible name, beautiful algorithm.")

s = slide("The travel analogy. You don't consult a street map of the whole country. "
          "Flight, train, walk. Long hops first, short hops last.")
head(s, "Flight → train → walk", "4 · hnsw")
panel(s, 0.85, 2.0, 11.6, 1.3)
tx(s, "You're in a small town in Japan and need a café in Osaka.\n"
      "You don't consult a street map of the entire country.",
   1.2, 2.28, 10.9, 0.9, 21, AMBER, True, PP_ALIGN.CENTER, spacing=1.25)
code(s, [
    "L2:   A ------------------------ Z      few nodes, long hops",
    "L1:   A ------- M -------------- Z",
    "L0:   A - c - f - M - p - t - w - Z     EVERY node, short hops",
], top=3.65, size=17)
tx(s, "If you know skip lists: it's a skip list where the base layer is a proximity graph.",
   0.85, 5.5, 11.6, 0.5, 18, MUTED, italic=True, align=PP_ALIGN.CENTER)
tx(s, "Roughly O(log N) hops.", 0.85, 6.1, 11.6, 0.6, 24, GREEN, True, PP_ALIGN.CENTER)

s = slide("The search, in one breath. Greedily walk to whichever neighbour is closer. When none "
          "improves, drop a layer. At layer 0 keep a candidate list of size efSearch.")
head(s, "The search, in one breath", "4 · hnsw")
bullets(s, [
    ("Start at the entry point in the top layer.", FG),
    ("Greedily walk to whichever neighbour is closer to the query.", FG),
    ("When no neighbour improves — drop down a layer. Repeat.", FG),
    ("At layer 0, keep a candidate list of size efSearch instead of just the best.", AMBER, True),
], top=2.35, gap=0.72)
panel(s, 0.85, 5.35, 11.6, 1.15)
tx(s, "Every node lives at layer 0. Upper layers are a routing shortcut, not a filter.",
   1.2, 5.62, 10.9, 0.6, 20, GREEN, True, PP_ALIGN.CENTER)

s = slide("Read it right to left, the way a search runs. Five nodes span the whole space up top; "
          "layer 0 has all 500. That's the flight, the train, and the walk.")
head(s, "The actual graph", "4 · hnsw")
image(s, "hnsw_layers.png", top=2.1, height=3.5)
tx(s, "Read right to left — the way a search runs. Sparse express lanes on top, every node at the bottom.",
   0.85, 5.9, 11.6, 0.5, 18, MUTED, align=PP_ALIGN.CENTER)
tx(s, f"levels: {N['levels_hist'][1]} nodes stop at layer 0  ·  {N['levels_hist'][2]} reach layer 1  ·  {N['levels_hist'][3]} reach layer 2",
   0.85, 6.4, 11.6, 0.5, 16, DIM, align=PP_ALIGN.CENTER)

hn = {r["ef"]: r for r in N["hnsw_sweep"]}
s = slide("Same index, same box, no rebuild — a per-request parameter. Most teams treat this as "
          "a config-file setting. It isn't.")
head(s, "One index, three products", "4 · hnsw")
table(s, ["efSearch", "ms/query", "recall@10", "plausible use case"],
      [(f"{ef}", f"{hn[ef]['ms']:.3f}", f"{hn[ef]['recall']:.3f}", uc)
       for ef, uc in [(10, "type-ahead"), (50, "standard search"), (100, ""),
                      (200, "RAG / agent context"), (400, "offline eval")]],
      [2.0, 2.3, 2.4, 5.0], top=2.3, size=17, hi=None, row_h=0.56)
tx(s, "Same index. Same box. No rebuild. It's a per-request parameter.",
   0.85, 5.6, 11.6, 0.6, 23, GREEN, True, PP_ALIGN.CENTER)
tx(s, "IVF's nprobe gives you the same freedom. Most teams leave this capability on the floor.",
   0.85, 6.25, 11.6, 0.5, 17, MUTED, align=PP_ALIGN.CENTER)

sc = N["scaling"]
s = slide("Not the latency at one size — the SHAPE of the curve as you grow. That's the whole pitch.")
head(s, "O(log N) vs O(N) — the shape is the argument", "4 · hnsw")
image(s, "scaling.png", top=1.95, height=3.9)
g_exact = sc[-1]["exact"] / sc[0]["exact"]
g_hnsw = sc[-1]["hnsw"] / sc[0]["hnsw"]
tx(s, f"Corpus grew {sc[-1]['n']/sc[0]['n']:.0f}×   →   exact got {g_exact:.1f}× slower,   HNSW got {g_hnsw:.1f}× slower.",
   0.85, 6.05, 11.6, 0.6, 21, AMBER, True, PP_ALIGN.CENTER)
tx(s, "And at 12,500 vectors HNSW is only 3× faster — the data agrees: at small N, don't bother.",
   0.85, 6.65, 11.6, 0.5, 16, MUTED, align=PP_ALIGN.CENTER)

s = slide("Three places, and they're serious. Memory is the headline — and it's your "
          "one-sentence comparison for the whole talk.")
head(s, "Where HNSW's value runs out", "4 · hnsw")
bullets(s, [
    (f"MEMORY — stores full vectors PLUS the graph: {N['hnsw_bytes']/1e6:.0f} MB vs IVF-PQ's {N['pq_bytes']/1e6:.1f} MB.", RED, True),
    (f"That's {N['hnsw_bytes']/N['pq_bytes']:.0f}× more memory.", MUTED),
    ("", FG),
    ("DELETES are genuinely awkward — you can't cut a node without risking disconnection.", FG),
    ("Tombstone and rebuild. For high-churn data that's a real operational cost.", MUTED),
    ("", FG),
    (f"BUILD TIME — {N['HNSW_BUILD_S']:.0f}s for 200k here. At 10M it's hours.", FG),
], top=2.15, gap=0.55)
panel(s, 0.85, 6.0, 11.6, 1.0)
tx(s, "HNSW buys speed with RAM.   IVF-PQ buys RAM with accuracy.",
   1.2, 6.27, 10.9, 0.6, 22, GREEN, True, PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════ 5 THE LEDGER
section("5", "The ledger", "Everything measured, on one problem, on one machine.",
        "Let me put all of it on one scoreboard.")

s = slide("Read the memory column and the recall column together. The bottom two rows are the "
          "same speed and nearly the same quality — and one uses 36x less memory.")
head(s, "The value ledger", "5 · scoreboard")
ivf16 = ivf[16]
table(s, ["step", "ms/query", "speedup", "memory", "recall", "what it bought"],
      [("Exact (numpy)", f"{N['EXACT_MS']:.3f}", "1×", f"{N['raw_bytes']/1e6:.0f} MB", "1.000",
        "correctness + ground truth"),
       ("+ IVF (nprobe=16)", f"{ivf16['ms']:.3f}", f"{N['EXACT_MS']/ivf16['ms']:.0f}×",
        f"{N['raw_bytes']/1e6:.0f} MB", f"{ivf16['recall']:.3f}", "32× less data scanned"),
       ("+ PQ (IVF-PQ)", f"{N['MS_IVFPQ']:.3f}", f"{N['EXACT_MS']/N['MS_IVFPQ']:.0f}×",
        f"{N['pq_bytes']/1e6:.1f} MB", f"{N['REC_IVFPQ']:.3f}", "32× less memory"),
       ("+ rerank (top-100)", f"{N['MS_RERANK']:.3f}", f"{N['EXACT_MS']/N['MS_RERANK']:.0f}×",
        f"{N['pq_bytes']/1e6:.1f} MB", f"{N['REC_RERANK']:.3f}", "recall bought back, ~free"),
       ("HNSW (ef=100)", f"{hn[100]['ms']:.3f}", f"{N['EXACT_MS']/hn[100]['ms']:.0f}×",
        f"{N['hnsw_bytes']/1e6:.0f} MB", f"{hn[100]['recall']:.3f}", "best recall/latency — for RAM")],
      [2.7, 1.5, 1.3, 1.5, 1.3, 3.3], top=2.25, size=14, hi=None, row_h=0.6)
tx(s, "The bottom two rows are the same speed and nearly the same quality —\nand one of them uses 36× less memory. That's your decision.",
   0.85, 5.85, 11.6, 1.0, 19, AMBER, True, spacing=1.25)

s = slide("The standard chart, built from our numbers. Up and to the right wins. Every technique "
          "in this talk was a way to push a point up and to the right.")
head(s, "The only chart that matters", "5 · scoreboard")
image(s, "recall_qps.png", top=1.9, height=4.5)
tx(s, "Every technique in this talk was a way to push a point up and to the right.",
   0.85, 6.5, 11.6, 0.5, 18, MUTED, align=PP_ALIGN.CENTER)

s = slide("Read straight down. And the one I'd lead with in a design review: if you're already "
          "on Postgres, pgvector gives you HNSW in the database you already run.")
head(s, "How to choose — read straight down", "5 · scoreboard")
y = 2.15
for cond, ans, col in [
        ("Under 100,000 vectors", "Brute force. Three lines of numpy. Don't add an index.", GREEN),
        ("Fits in RAM, want best quality/latency", "HNSW. The default in every major system for a reason.", BLUE),
        ("RAM is your constraint, or past ~100M", "IVF-PQ, with reranking. Non-negotiable on the reranking.", AMBER),
        ("Always", "Keep exact search around as the ruler.", MUTED)]:
    panel(s, 0.85, y, 11.6, 0.95)
    tx(s, cond, 1.15, y + 0.28, 4.6, 0.5, 17, col, True)
    tx(s, ans, 6.0, y + 0.3, 6.2, 0.5, 16, FG)
    y += 1.08
tx(s, "Two indexes. Pick by which resource you're short of.\n"
      "And if you're already on Postgres — pgvector gives you HNSW in the database you already run.",
   0.85, 6.5, 11.6, 0.9, 17, MUTED, spacing=1.25)

# ═══════════════════════════════════════════════════════════ 6 FILTERING
section("6", "The gotcha", "Everything so far assumed the query is JUST a vector.\nIt never is.",
        "Four minutes on the thing that will actually bite you.")

s = slide("Real queries are never just a vector. Let's break it on purpose.")
head(s, "Real queries are never just a vector", "6 · filtering")
panel(s, 0.85, 2.05, 11.6, 1.15)
tx(s, "Nearest neighbours WHERE tenant_id = 47 AND status = 'open' AND created_at > monday",
   1.2, 2.4, 10.9, 0.5, 18, AMBER, True, PP_ALIGN.CENTER, font=MONO)
tx(s, "The naive approach — POST-FILTERING: run normal ANN search, then drop the results\nthat fail the predicate. Watch what happens when the predicate is selective.",
   0.85, 3.5, 11.6, 1.0, 19, FG, spacing=1.3)
stat(s, f"{N['kept_mean']:.2f} / 10", "results returned, on average",
     f"{N['kept_zero_pct']:.0f}% of queries returned ZERO results — with {N['n_priority']:,} matching tickets in the corpus",
     RED, top=4.35, vsize=60)

s = slide("That is a broken product feature, not a tuning problem. The fix is to push the "
          "predicate INTO the search so the index keeps walking until it has enough matches.")
head(s, "Push the predicate INTO the search", "6 · filtering")
table(s, ["strategy", "avg results out of 10"],
      [("post-filter (search, then drop)", f"{N['kept_mean']:.2f}"),
       ("IVF + IDSelector (pushed down)", f"{N['ivf_valid']:.2f}"),
       ("HNSW + IDSelector (pushed down)", f"{N['hnsw_valid']:.2f}")],
      [7.4, 4.2], top=2.4, size=18, hi=0, hi_color=RED, row_h=0.62)
tx(s, "Same index. Same query. Same filter.\n0.18 versus 10 — entirely determined by whether the predicate went in or after.",
   0.85, 4.9, 11.6, 1.0, 20, AMBER, True, spacing=1.25)
tx(s, "Pre-filtering wins when the filter is very selective · filtered graph traversal degrades when matches are scattered\n"
      "Every vector database handles this differently — and this, more than raw QPS, should decide which one you pick.",
   0.85, 6.15, 11.6, 1.0, 15, MUTED, spacing=1.35)

# ═══════════════════════════════════════════════════════════════════ 7 CLOSE
s = slide("Four techniques. Three things I'd like you to actually leave with.")
head(s, "Three things to leave with")
y = 2.2
for i, (t, d, col) in enumerate([
        ("Check your row count first.",
         "Under 100k vectors the answer is three lines of numpy. Not a vector database.", GREEN),
        ("If you use PQ, rerank.",
         f"{N['REC_IVFPQ']:.3f} → {N['REC_RERANK']:.3f} for {(N['MS_RERANK']-N['MS_IVFPQ'])*1000:.0f} microseconds. "
         "The cheapest quality win in the field, and routinely missing.", AMBER),
        ("Recall degrades silently.",
         "No exception, no alert, no 500 — just quietly worse results. Ship the ground-truth "
         "job WITH the index, not after the incident.", RED)]):
    panel(s, 0.85, y, 11.6, 1.35)
    tx(s, f"{i+1}", 1.15, y + 0.35, 0.6, 0.6, 30, col, True)
    tx(s, t, 1.95, y + 0.24, 10.2, 0.5, 21, col, True)
    tx(s, d, 1.95, y + 0.72, 10.2, 0.6, 16, MUTED)
    y += 1.5

s = slide("The notebook is the deliverable. Every number here comes out of it. Questions.")
tx(s, "Questions", 0.9, 2.5, 11.6, 1.2, 56, FG, True)
rule(s, 0.92, 3.85, 2.6, AMBER)
tx(s, "Everything in this deck is measured by:", 0.9, 4.3, 11.6, 0.5, 19, MUTED)
tx(s, "notebooks/talk-vector-search.ipynb", 0.9, 4.75, 11.6, 0.5, 22, BLUE, True, font=MONO)
tx(s, "Runs in ~5 minutes on a free Colab CPU. Change N_TICKETS at the top and every\n"
      "table and chart recomputes. There's a “Things to try” section with seven experiments.",
   0.9, 5.5, 11.6, 1.0, 16, DIM, spacing=1.3)

# ═════════════════════════════════════════════════════════════════════ WRITE
out = os.path.join(HERE, "vector-search-talk.pptx")
prs.save(out)
print(f"wrote {out}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
print(f"with speaker notes: {sum(1 for n in deck_notes if n)}")
